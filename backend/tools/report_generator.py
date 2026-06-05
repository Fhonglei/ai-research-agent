import os
import re
import markdown as md_lib
from pathlib import Path
from typing import Optional
from models.schemas import ResearchReport, ResearchTask, Source
from utils.logger import logger


REPORTS_DIR = Path(__file__).resolve().parent.parent / "reports"

# Professional CSS stylesheet for PDF rendering
PDF_CSS = """
@page {
    size: A4;
    margin: 2.5cm 2cm 2.5cm 2cm;
    @top-center {
        content: element(header);
        font-family: 'Helvetica', 'Arial', sans-serif;
        font-size: 9pt;
        color: #888;
    }
    @bottom-center {
        content: counter(page);
        font-family: 'Helvetica', 'Arial', sans-serif;
        font-size: 9pt;
        color: #888;
    }
}

body {
    font-family: 'Helvetica', 'Arial', sans-serif;
    font-size: 11pt;
    line-height: 1.6;
    color: #1a1a1a;
    counter-reset: section;
}

h1 {
    font-size: 24pt;
    color: #1a3a5c;
    border-bottom: 3px solid #1a3a5c;
    padding-bottom: 10px;
    margin-top: 0;
    page-break-before: avoid;
}

h2 {
    font-size: 16pt;
    color: #2b5f8e;
    border-bottom: 1px solid #c0d4e8;
    padding-bottom: 6px;
    margin-top: 28px;
    page-break-before: avoid;
}

h3 {
    font-size: 13pt;
    color: #3a7cc3;
    margin-top: 20px;
    page-break-before: avoid;
}

h4 {
    font-size: 11pt;
    color: #555;
    margin-top: 16px;
}

p {
    margin: 8px 0;
    text-align: justify;
}

ul, ol {
    margin: 8px 0;
    padding-left: 24px;
}

li {
    margin: 4px 0;
}

blockquote {
    border-left: 4px solid #c0d4e8;
    margin: 12px 0;
    padding: 8px 16px;
    background: #f7f9fc;
    font-style: italic;
}

a {
    color: #2b5f8e;
    text-decoration: none;
    word-wrap: break-word;
}

code {
    font-family: 'Consolas', 'Courier New', monospace;
    font-size: 9pt;
    background: #f0f0f0;
    padding: 1px 4px;
    border-radius: 2px;
}

pre {
    background: #f7f9fc;
    border: 1px solid #ddd;
    padding: 12px;
    font-size: 9pt;
    overflow-x: auto;
    border-radius: 4px;
}

pre code {
    background: none;
    padding: 0;
}

table {
    width: 100%;
    border-collapse: collapse;
    margin: 12px 0;
    font-size: 10pt;
}

th, td {
    border: 1px solid #ddd;
    padding: 8px 10px;
    text-align: left;
}

th {
    background: #1a3a5c;
    color: white;
    font-weight: 600;
}

tr:nth-child(even) {
    background: #f7f9fc;
}

hr {
    border: none;
    border-top: 1px solid #ddd;
    margin: 20px 0;
}

.toc {
    background: #f7f9fc;
    border: 1px solid #c0d4e8;
    border-radius: 6px;
    padding: 16px 20px;
    margin: 20px 0;
    page-break-before: avoid;
}

.toc h2 {
    margin-top: 0;
    border-bottom: none;
    font-size: 14pt;
}

.toc ul {
    list-style: none;
    padding-left: 0;
}

.toc li {
    margin: 6px 0;
}

.executive-summary {
    background: #eaf2f8;
    border-left: 5px solid #2b5f8e;
    padding: 16px 20px;
    margin: 16px 0;
    border-radius: 0 6px 6px 0;
}

.key-takeaway {
    background: #fef9e7;
    border-left: 5px solid #f0b429;
    padding: 12px 16px;
    margin: 8px 0;
    border-radius: 0 6px 6px 0;
}

.reference {
    font-size: 9pt;
    padding: 2px 0;
    word-break: break-all;
}
"""


def generate_markdown(report: ResearchReport) -> str:
    """
    Generate a professional markdown report from a ResearchReport.

    Includes a table of contents, executive summary (derived from all task
    summaries), individual subtopic sections with source citations, key
    takeaways, and a combined references section.

    Args:
        report: The ResearchReport containing tasks and metadata.

    Returns:
        A complete markdown document string.
    """
    if not report.topic:
        logger.warning("generate_markdown called with empty topic")
        return "# Research Report\n\n*No topic provided.*"

    lines: list[str] = []

    # Title
    lines.append(f"# Research Report: {report.topic}")
    lines.append("")
    lines.append(f"**Date:** {report.created_at[:10] if report.created_at else 'N/A'}")
    lines.append(f"**Depth:** {report.depth or 'standard'} | "
                 f"**Sources consulted:** {sum(len(t.sources) for t in report.tasks)}")
    lines.append("")

    # Table of Contents
    lines.append("## Table of Contents")
    lines.append("")
    lines.append("1. [Executive Summary](#executive-summary)")
    for i, subtopic in enumerate(report.subtopics, start=2):
        anchor = _slugify(subtopic)
        lines.append(f"{i}. [{subtopic}](#{anchor})")
    lines.append(f"{len(report.subtopics) + 2}. [Key Takeaways](#key-takeaways)")
    lines.append(f"{len(report.subtopics) + 3}. [References](#references)")
    lines.append("")

    # Executive Summary — weave all task summaries together
    lines.append("## Executive Summary")
    lines.append("")
    if report.tasks:
        combined = " ".join(t.summary for t in report.tasks if t.summary)
        if combined:
            lines.append(combined)
        else:
            lines.append("*Analysis in progress...*")
    else:
        lines.append("*No research tasks were completed.*")
    lines.append("")

    # Subtopic sections
    for task in report.tasks:
        if not task.subtopic:
            continue
        lines.append(f"## {task.subtopic}")
        lines.append("")

        if task.summary:
            lines.append(task.summary)
        else:
            lines.append("*Research pending...*")

        if task.sources:
            lines.append("")
            lines.append("### Sources")
            lines.append("")
            for j, src in enumerate(task.sources, start=1):
                title_display = src.title or src.url
                lines.append(f"{j}. **[{title_display}]({src.url})**")
                if src.snippet:
                    lines.append(f"   - {src.snippet}")
            lines.append("")

    # Key Takeaways
    lines.append("## Key Takeaways")
    lines.append("")
    if report.tasks:
        for i, task in enumerate(report.tasks, start=1):
            if task.summary:
                # Extract first sentence or key phrase
                first_sentence = task.summary.split(".")[0]
                if first_sentence:
                    lines.append(f"- **{task.subtopic}:** {first_sentence.strip()}.")
    else:
        lines.append("- *No takeaways available.*")
    lines.append("")

    # References
    lines.append("## References")
    lines.append("")
    ref_index = 1
    for task in report.tasks:
        for src in task.sources:
            if src.url:
                title = src.title or "Untitled"
                lines.append(f"{ref_index}. [{title}]({src.url})")
                ref_index += 1
    if ref_index == 1:
        lines.append("*No references collected.*")
    lines.append("")

    content = "\n".join(lines)
    logger.info(f"Generated markdown report: {len(content)} characters")
    return content


def markdown_to_pdf(markdown_content: str, output_path: str) -> str:
    """
    Convert markdown content to a PDF file using WeasyPrint.

    The markdown is first rendered to HTML, then styled with a professional
    CSS theme and converted to PDF.

    Args:
        markdown_content: The markdown string to convert.
        output_path: Filesystem path to write the PDF to.

    Returns:
        The output_path on success, empty string on failure.
    """
    if not markdown_content:
        logger.warning("markdown_to_pdf called with empty content")
        return ""

    try:
        # Ensure output directory exists
        out_path = Path(output_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)

        # Convert markdown to HTML
        md = md_lib.Markdown(extensions=["extra", "codehilite", "toc", "sane_lists"])
        html_body = md.convert(markdown_content)

        # Build a complete HTML document with CSS
        html_document = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<style>
{PDF_CSS}
</style>
</head>
<body>
{html_body}
</body>
</html>"""

        # Convert to PDF using WeasyPrint
        from weasyprint import HTML
        HTML(string=html_document).write_pdf(str(output_path))

        logger.info(f"PDF generated at: {output_path}")
        return str(output_path)

    except ImportError:
        logger.error("WeasyPrint is not installed. Install with: pip install weasyprint")
        return ""
    except Exception as e:
        logger.error(f"Failed to generate PDF: {e}")
        return ""


def markdown_to_pptx(markdown_content: str, output_path: str) -> str:
    """
    Convert markdown content to a PowerPoint (PPTX) presentation.

    Creates a title slide followed by one slide per major section (H2 headings).
    Each section slide contains bullet points derived from the text content.

    Args:
        markdown_content: The markdown string to convert.
        output_path: Filesystem path to write the PPTX to.

    Returns:
        The output_path on success, empty string on failure.
    """
    if not markdown_content:
        logger.warning("markdown_to_pptx called with empty content")
        return ""

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Parse markdown into sections
        sections = _parse_markdown_sections(markdown_content)
        title_text = sections.pop(0) if sections else "Research Report"

        # ── Title Slide ──
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)

        # Title text
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(10.333)
        height = Inches(2.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        top2 = Inches(4.5)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.0))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "AI-Generated Research Report"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xC0, 0xD4, 0xE8)
        p2.alignment = PP_ALIGN.LEFT

        # ── Section Slides ──
        for section_title, bullets in sections:
            slide = prs.slides.add_slide(slide_layout)

            # Header bar
            header_shape = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
            header_shape.line.fill.background()

            tf_header = header_shape.text_frame
            tf_header.word_wrap = True
            p_header = tf_header.paragraphs[0]
            p_header.text = section_title
            p_header.font.size = Pt(28)
            p_header.font.bold = True
            p_header.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p_header.alignment = PP_ALIGN.LEFT

            # Bullet points
            left_margin = Inches(1.0)
            top_margin = Inches(1.8)
            content_width = Inches(11.333)
            content_height = Inches(5.0)
            txBox = slide.shapes.add_textbox(left_margin, top_margin, content_width, content_height)
            tf = txBox.text_frame
            tf.word_wrap = True

            if bullets:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet.strip()
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                    p.space_after = Pt(8)
                    p.level = 0
            else:
                p = tf.paragraphs[0]
                p.text = "(No content available for this section)"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

        # Save
        out_path = Path(output_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        logger.info(f"PPTX generated at: {output_path}")
        return str(output_path)

    except ImportError as e:
        logger.error(f"Missing dependency for PPTX generation: {e}")
        return ""
    except Exception as e:
        logger.error(f"Failed to generate PPTX: {e}")
        return ""


def _parse_markdown_sections(markdown_content: str) -> list[tuple[str, list[str]]]:
    """
    Parse markdown content into sections based on H1/H2 headings.

    Args:
        markdown_content: Raw markdown string.

    Returns:
        List of (section_title, list_of_bullet_lines) tuples.
    """
    lines = markdown_content.split("\n")
    sections: list[tuple[str, list[str]]] = []
    current_title = "Research Report"
    current_bullets: list[str] = []
    seen_headings: set[str] = set()

    for line in lines:
        stripped = line.strip()

        # H1 — main title
        if stripped.startswith("# ") and not stripped.startswith("## "):
            current_title = stripped[2:].strip()
            if not sections:
                sections.append((current_title, []))
                current_bullets = []
            continue

        # H2 — new section
        if stripped.startswith("## "):
            heading = stripped[3:].strip()
            if heading.lower() in ("table of contents", "references"):
                continue
            # Save previous section
            if current_bullets:
                for i in range(len(sections)):
                    if sections[i][0] == current_title:
                        sections[i] = (current_title, list(current_bullets))
                if not any(s[0] == current_title for s in sections):
                    sections.append((current_title, list(current_bullets)))
            current_title = heading
            current_bullets = []
            continue

        # Bullet points
        if stripped.startswith("- ") or stripped.startswith("* "):
            bullet_text = stripped[2:].strip()
            if bullet_text and bullet_text not in current_bullets:
                current_bullets.append(bullet_text)
            continue

        # Regular text — append as bullet if substantial
        if stripped and not stripped.startswith(">") and not stripped.startswith("```"):
            # Skip text that is just a standalone reference or URL
            if stripped.startswith("[") and stripped.endswith(")"):
                continue
            if stripped.startswith("http"):
                continue
            if len(stripped) > 20 and stripped not in current_bullets:
                current_bullets.append(stripped)

    # Save final section
    if current_bullets:
        for i in range(len(sections)):
            if sections[i][0] == current_title:
                sections[i] = (current_title, list(current_bullets))
        if not any(s[0] == current_title for s in sections):
            sections.append((current_title, list(current_bullets)))

    # Deduplicate bullets within each section
    result = []
    for title, bullets in sections:
        seen = set()
        deduped = []
        for b in bullets:
            if b not in seen:
                seen.add(b)
                deduped.append(b)
        result.append((title, deduped))

    return result if result else [("Research Report", [])]


def _slugify(text: str) -> str:
    """
    Convert a section title to a markdown-anchor-compatible slug.

    Args:
        text: The heading text.

    Returns:
        Lowercase, hyphen-separated slug.
    """
    slug = text.lower().strip()
    slug = re.sub(r"[^\w\s-]", "", slug)
    slug = re.sub(r"[\s_]+", "-", slug)
    slug = re.sub(r"-+", "-", slug)
    return slug.strip("-")
